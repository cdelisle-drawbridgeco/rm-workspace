"""
Generate the Drawbridge Business Outcome Framework PowerPoint deck.
Run: python3 scripts/generate-outcome-deck.py
Output: Drawbridge_Outcome_Framework.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# Brand colors
DB_DARK = RGBColor(0x11, 0x18, 0x21)
DB_AQUA = RGBColor(0x86, 0xC8, 0xBC)
DB_AQUA_DARK = RGBColor(0x48, 0x9E, 0x8E)
DB_ORANGE = RGBColor(0xFF, 0x69, 0x00)
DB_GRAY_LIGHT = RGBColor(0xF6, 0xF6, 0xF7)
DB_GRAY = RGBColor(0xD8, 0xD9, 0xDB)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK_TEXT = RGBColor(0x22, 0x22, 0x22)
MID_TEXT = RGBColor(0x55, 0x55, 0x55)
LIGHT_TEXT = RGBColor(0x88, 0x88, 0x88)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


# ── Helpers ──────────────────────────────────────────────────────────

def add_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, fill_color, border_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1.5)
    else:
        shape.line.fill.background()
    # Smaller corner radius
    shape.adjustments[0] = 0.05
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=14,
                 color=DARK_TEXT, bold=False, alignment=PP_ALIGN.LEFT,
                 font_name='Calibri'):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_para(text_frame, text, font_size=14, color=DARK_TEXT, bold=False,
             space_before=Pt(4), space_after=Pt(2), alignment=PP_ALIGN.LEFT,
             font_name='Calibri'):
    p = text_frame.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    p.space_before = space_before
    p.space_after = space_after
    return p


def add_bullet(text_frame, text, font_size=13, color=DARK_TEXT, level=0,
               bold_prefix=None):
    p = text_frame.add_paragraph()
    p.level = level
    p.space_before = Pt(3)
    p.space_after = Pt(2)
    if bold_prefix:
        run_b = p.add_run()
        run_b.text = bold_prefix + " "
        run_b.font.size = Pt(font_size)
        run_b.font.color.rgb = color
        run_b.font.bold = True
        run_b.font.name = 'Calibri'
        run = p.add_run()
        run.text = text
        run.font.size = Pt(font_size)
        run.font.color.rgb = color
        run.font.bold = False
        run.font.name = 'Calibri'
    else:
        p.text = text
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = 'Calibri'
    return p


def dark_header_bar(slide):
    """Top dark bar with Drawbridge branding."""
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.9), DB_DARK)
    add_text_box(slide, Inches(0.6), Inches(0.18), Inches(3), Inches(0.55),
                 "DRAWBRIDGE", font_size=22, color=WHITE, bold=True)


def section_title_slide(title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    add_bg(slide, DB_DARK)
    # Accent line
    add_rect(slide, Inches(0.8), Inches(2.8), Inches(1.2), Inches(0.06), DB_AQUA)
    add_text_box(slide, Inches(0.8), Inches(3.0), Inches(10), Inches(1.2),
                 title, font_size=40, color=WHITE, bold=True)
    add_text_box(slide, Inches(0.8), Inches(4.2), Inches(10), Inches(0.8),
                 subtitle, font_size=18, color=DB_AQUA)
    # Logo text
    add_text_box(slide, Inches(0.8), Inches(6.4), Inches(3), Inches(0.5),
                 "DRAWBRIDGE", font_size=14, color=RGBColor(0x55, 0x66, 0x70), bold=True)
    return slide


# ── SLIDE 1: Title ──────────────────────────────────────────────────

slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DB_DARK)

add_text_box(slide, Inches(0.8), Inches(1.5), Inches(5), Inches(0.5),
             "DRAWBRIDGE", font_size=16, color=DB_AQUA, bold=True)

# Accent line
add_rect(slide, Inches(0.8), Inches(2.2), Inches(1.5), Inches(0.06), DB_ORANGE)

add_text_box(slide, Inches(0.8), Inches(2.5), Inches(11), Inches(1.5),
             "Business Outcome\nFramework", font_size=48, color=WHITE, bold=True)

add_text_box(slide, Inches(0.8), Inches(4.5), Inches(10), Inches(0.8),
             "Moving from features & deliverables to strategic business outcomes\nthat drive client conversations, renewals, and expansion.",
             font_size=18, color=DB_GRAY)

add_text_box(slide, Inches(0.8), Inches(6.2), Inches(5), Inches(0.5),
             "INTERNAL USE ONLY  |  CONFIDENTIAL", font_size=11, color=LIGHT_TEXT)


# ── SLIDE 2: Why Outcomes ───────────────────────────────────────────

slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
dark_header_bar(slide)

add_text_box(slide, Inches(0.6), Inches(1.2), Inches(12), Inches(0.6),
             "Why Outcomes, Not Features", font_size=32, color=DB_DARK, bold=True)

# Left: the problem
box_l = add_rect(slide, Inches(0.6), Inches(2.0), Inches(5.8), Inches(4.8),
                 RGBColor(0xFF, 0xF0, 0xF0), RGBColor(0xE8, 0xC0, 0xC0))
tf = box_l.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Feature-Led (Vendor Language)"
p.font.size = Pt(16)
p.font.color.rgb = RGBColor(0xCC, 0x33, 0x33)
p.font.bold = True
p.font.name = 'Calibri'

add_para(tf, '"This year we completed 4 assessments, updated 12 policies, ran 2 tabletop exercises, and responded to 47 DDQs. Here\'s the renewal for the same scope."',
         font_size=13, color=MID_TEXT, space_before=Pt(16))
add_para(tf, "", font_size=8)
add_bullet(tf, "Sounds like a task report", font_size=12, color=MID_TEXT)
add_bullet(tf, "Client hears: cost, not value", font_size=12, color=MID_TEXT)
add_bullet(tf, "Easy to compare on price", font_size=12, color=MID_TEXT)
add_bullet(tf, "Positions us as a vendor", font_size=12, color=MID_TEXT)

# Right: the solution
box_r = add_rect(slide, Inches(6.8), Inches(2.0), Inches(5.8), Inches(4.8),
                 RGBColor(0xF0, 0xFF, 0xF0), RGBColor(0xC0, 0xE8, 0xC0))
tf = box_r.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Outcome-Led (Client Language)"
p.font.size = Pt(16)
p.font.color.rgb = RGBColor(0x33, 0x88, 0x33)
p.font.bold = True
p.font.name = 'Calibri'

add_para(tf, '"This year, we helped you raise your new fund with confidence — your ODD reviews went smoothly and DDQ turnaround dropped from 3 weeks to 5 days. We kept you exam-ready year-round. And we did it without you hiring a single security person."',
         font_size=13, color=MID_TEXT, space_before=Pt(16))
add_para(tf, "", font_size=8)
add_bullet(tf, "Sounds like a strategic partner", font_size=12, color=MID_TEXT)
add_bullet(tf, "Client hears: impact and ROI", font_size=12, color=MID_TEXT)
add_bullet(tf, "Hard to replace — we solve problems", font_size=12, color=MID_TEXT)
add_bullet(tf, "Positions us as infrastructure", font_size=12, color=MID_TEXT)


# ── SLIDE 3: The Six Outcomes Overview ─────────────────────────────

slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
dark_header_bar(slide)

add_text_box(slide, Inches(0.6), Inches(1.2), Inches(12), Inches(0.6),
             "The Six Business Outcomes", font_size=32, color=DB_DARK, bold=True)

add_text_box(slide, Inches(0.6), Inches(1.8), Inches(11), Inches(0.5),
             "Every engagement should map to one or more of these outcomes.",
             font_size=14, color=MID_TEXT)

outcomes = [
    ("1", "Raise & Retain Capital\nwith Confidence",
     "Investor trust through\noperational excellence"),
    ("2", "Navigate Regulatory\nComplexity",
     "Stay ahead of evolving\nregulations"),
    ("3", "Protect the Firm\nfrom Disruption",
     "Be ready for the incident\nthat hasn't happened yet"),
    ("4", "Make Security a\nCompetitive Advantage",
     "Turn cost center into\ndifferentiator"),
    ("5", "Scale Without\nAdding Headcount",
     "Enterprise security without\nenterprise team"),
    ("6", "Get AI-Ready,\nStay AI-Safe",
     "Adopt AI with confidence —\ntraining, assessments & governance"),
]

card_w = Inches(1.9)
card_h = Inches(4.2)
start_x = Inches(0.6)
gap = Inches(0.18)

for i, (num, title, subtitle) in enumerate(outcomes):
    x = start_x + i * (card_w + gap)
    y = Inches(2.5)

    # Card background
    card = add_rect(slide, x, y, card_w, card_h, DB_GRAY_LIGHT, DB_AQUA)

    # Number circle
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.6), y + Inches(0.3),
                                     Inches(0.65), Inches(0.65))
    circle.fill.solid()
    circle.fill.fore_color.rgb = DB_AQUA
    circle.line.fill.background()
    tf = circle.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = num
    p.font.size = Pt(22)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.font.name = 'Calibri'
    p.alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].space_before = Pt(0)

    # Title
    add_text_box(slide, x + Inches(0.1), y + Inches(1.2), card_w - Inches(0.2), Inches(1.2),
                 title, font_size=14, color=DB_DARK, bold=True, alignment=PP_ALIGN.CENTER)

    # Subtitle
    add_text_box(slide, x + Inches(0.1), y + Inches(2.6), card_w - Inches(0.2), Inches(1.0),
                 subtitle, font_size=10, color=MID_TEXT, alignment=PP_ALIGN.CENTER)


# ── SLIDES 4-8: Individual Outcome Deep Dives ──────────────────────

outcome_details = [
    {
        "num": "1",
        "title": "Raise & Retain Capital with Confidence",
        "subtitle": "Investor trust is earned through operational excellence — not just performance.",
        "why": (
            "Allocators commit capital to managers they trust — and trust increasingly requires "
            "demonstrating operational and cybersecurity maturity. A single failed DDQ response, "
            "an incomplete ODD review, or a data breach headline can cost a firm hundreds of "
            "millions in commitments."
        ),
        "impacts": [
            "Faster DDQ turnaround with higher-quality, consistent responses",
            "Smoother LP due diligence cycles with fewer follow-up requests",
            "Confidence to pursue new allocator relationships",
            "Reduced risk of capital loss from operational deficiency concerns",
        ],
        "client_says": [
            '"Our LPs keep asking about cybersecurity"',
            '"DDQ season is killing us"',
            '"We\'re raising a new fund"',
            '"Our consultant flagged our security program"',
            '"We need to look institutional"',
        ],
    },
    {
        "num": "2",
        "title": "Navigate Regulatory Complexity",
        "subtitle": "Stay ahead of evolving regulations without building a compliance department.",
        "why": (
            "The regulatory landscape is expanding and accelerating — Reg S-P amendments, SEC exam "
            "priorities, DORA in Europe, state privacy laws. Keeping up is a full-time job that "
            "most alternative managers don't have the headcount for. The cost of non-compliance "
            "isn't just fines — it's reputational damage and LP concern."
        ),
        "impacts": [
            "Exam-ready posture maintained year-round",
            "Proactive compliance with new rules before enforcement begins",
            "Reduced legal and consulting spend on regulatory interpretation",
            "Confidence the firm won't be caught off guard by a rule change",
        ],
        "client_says": [
            '"What do the new SEC rules mean for us?"',
            '"We got an exam notification"',
            '"Our outside counsel says we need to update policies"',
            '"We have European investors now"',
            '"We just want to stay out of trouble"',
        ],
    },
    {
        "num": "3",
        "title": "Protect the Firm from Disruption",
        "subtitle": "Be ready for the incident that hasn't happened yet — and recover fast when it does.",
        "why": (
            "Cyber incidents aren't theoretical — they're statistical certainties. Ransomware can "
            "shut down operations for weeks. A BEC scam can move millions in minutes. A data breach "
            "can trigger regulatory action, LP redemptions, and lasting reputational harm. The firms "
            "that survive these events are the ones that prepared."
        ),
        "impacts": [
            "Reduced likelihood and blast radius of security incidents",
            "Faster detection and response when incidents occur",
            "Tested recovery procedures that work under pressure",
            "Lower financial and reputational cost of incidents",
        ],
        "client_says": [
            '"What happens if we get hit?"',
            '"Our peer firm just had a breach"',
            '"We need someone to call at 2am"',
            '"We\'re worried about ransomware"',
            '"Do we have a plan if something happens?"',
        ],
    },
    {
        "num": "4",
        "title": "Make Security a Competitive Advantage",
        "subtitle": "Turn a cost center into a differentiator that wins capital, talent, and trust.",
        "why": (
            "Most firms think of cybersecurity as a cost. The best-run firms think of it as a "
            "competitive weapon. In a market where allocators choose between managers with similar "
            "strategies and returns, the firm with the best operational infrastructure wins. A mature "
            "security program signals professionalism and institutional quality."
        ),
        "impacts": [
            "Differentiation in competitive fundraising situations",
            "Faster LP onboarding with pre-built compliance evidence",
            "Ability to win larger, more sophisticated institutional allocators",
            "Talent attraction — top candidates want well-run firms",
        ],
        "client_says": [
            '"We want to look best-in-class"',
            '"How do we compare to our peers?"',
            '"We\'re going after bigger LPs"',
            '"Our competitor just got SOC 2 certified"',
            '"We need to level up our operations"',
        ],
    },
    {
        "num": "5",
        "title": "Scale Without Adding Headcount",
        "subtitle": "Grow the firm, launch funds, expand operations — without building an internal security team.",
        "why": (
            "Alternative managers are lean by design — every headcount dollar that goes to "
            "non-investment functions is a dollar not generating alpha. But as firms grow, security "
            "demands grow too. Hiring a CISO ($300K-$500K+) plus a security team is prohibitive. "
            "They need enterprise-grade security that scales without the enterprise headcount."
        ),
        "impacts": [
            "Full security program at a fraction of in-house cost",
            "Instant access to specialized expertise on demand",
            "Security that grows with new funds, offices, and portfolio companies",
            "Free up COO / CCO time to focus on running the business",
        ],
        "client_says": [
            '"We don\'t have a CISO and don\'t want to hire one"',
            '"Our COO is stretched thin"',
            '"We just launched a new fund"',
            '"We\'re opening a London office"',
            '"We need this handled — I don\'t have the bandwidth"',
        ],
    },
    {
        "num": "6",
        "title": "Get AI-Ready, Stay AI-Safe",
        "subtitle": "Adopt AI with confidence — with the training, assessments, and governance to minimize risk.",
        "why": (
            "Every firm is being asked about AI — by LPs, by regulators, and by their own teams. "
            "Employees are already using AI tools, often without oversight. The SEC is paying "
            "attention to AI governance, and allocators are starting to ask 'what's your AI policy?' "
            "in DDQs. Firms that wait will find themselves reacting to incidents, regulatory gaps, "
            "and reputational risk. The firms that get ahead will have clear policies, trained teams, "
            "and a framework for evaluating and governing AI tools."
        ),
        "impacts": [
            "Employee AI training that reduces risk of data leakage and misuse",
            "LLM and AI tool assessments that identify exposure before it becomes a problem",
            "Clear, documented AI acceptable use policies that satisfy regulators and LPs",
            "Governance framework for evaluating, approving, and monitoring AI tools",
            "Confidence to adopt AI where it adds value — without creating new attack surfaces",
        ],
        "client_says": [
            '"Our team is using ChatGPT — should we be worried?"',
            '"LPs are asking about our AI policy"',
            '"We don\'t know what tools people are using"',
            '"The SEC mentioned AI in their exam priorities"',
            '"We want to use AI but don\'t know where to start"',
        ],
    },
]

for od in outcome_details:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, WHITE)
    dark_header_bar(slide)

    # Outcome number badge
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.6), Inches(1.15),
                                     Inches(0.6), Inches(0.6))
    circle.fill.solid()
    circle.fill.fore_color.rgb = DB_AQUA
    circle.line.fill.background()
    tf = circle.text_frame
    p = tf.paragraphs[0]
    p.text = od["num"]
    p.font.size = Pt(20)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.font.name = 'Calibri'
    p.alignment = PP_ALIGN.CENTER

    # Title
    add_text_box(slide, Inches(1.4), Inches(1.1), Inches(10), Inches(0.45),
                 od["title"], font_size=28, color=DB_DARK, bold=True)
    add_text_box(slide, Inches(1.4), Inches(1.55), Inches(10), Inches(0.35),
                 od["subtitle"], font_size=13, color=MID_TEXT)

    # Left column: Why It Matters + Business Impact
    # Why box
    why_box = add_rect(slide, Inches(0.6), Inches(2.2), Inches(7.4), Inches(2.0),
                       DB_GRAY_LIGHT)
    tf = why_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "WHY IT MATTERS"
    p.font.size = Pt(10)
    p.font.color.rgb = LIGHT_TEXT
    p.font.bold = True
    p.font.name = 'Calibri'
    add_para(tf, od["why"], font_size=12, color=DARK_TEXT, space_before=Pt(8))

    # Impact box
    impact_box = add_rect(slide, Inches(0.6), Inches(4.4), Inches(7.4), Inches(2.6),
                          DB_GRAY_LIGHT)
    tf = impact_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "THE BUSINESS IMPACT"
    p.font.size = Pt(10)
    p.font.color.rgb = LIGHT_TEXT
    p.font.bold = True
    p.font.name = 'Calibri'
    for imp in od["impacts"]:
        add_bullet(tf, imp, font_size=12, color=DARK_TEXT)

    # Right column: How Your Client Says It
    says_box = add_rect(slide, Inches(8.3), Inches(2.2), Inches(4.4), Inches(4.8),
                        RGBColor(0xE8, 0xF4, 0xF1), DB_AQUA)
    tf = says_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "HOW YOUR CLIENT SAYS IT"
    p.font.size = Pt(10)
    p.font.color.rgb = DB_AQUA_DARK
    p.font.bold = True
    p.font.name = 'Calibri'
    add_para(tf, "Listen for these phrases — they're buying signals:", font_size=11,
             color=MID_TEXT, space_before=Pt(10))
    for phrase in od["client_says"]:
        add_para(tf, phrase, font_size=12, color=DB_DARK, space_before=Pt(10))


# ── SLIDE 9: Client-Type Mapping ────────────────────────────────────

slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
dark_header_bar(slide)

add_text_box(slide, Inches(0.6), Inches(1.2), Inches(12), Inches(0.6),
             "Mapping Outcomes to Client Types", font_size=32, color=DB_DARK, bold=True)

add_text_box(slide, Inches(0.6), Inches(1.8), Inches(11), Inches(0.5),
             "Not every outcome matters equally to every client. Lead with what resonates most.",
             font_size=14, color=MID_TEXT)

# Table
rows = [
    ("Hedge Fund (fundraising)", "1. Capital Confidence", "2. Regulatory + 4. Competitive"),
    ("Hedge Fund (steady state)", "2. Regulatory Complexity", "3. Protection + 6. AI Readiness"),
    ("PE Firm (GP level)", "1. Capital Confidence", "5. Scale + 6. AI Readiness (portco)"),
    ("PE Portfolio Company", "3. Protect from Disruption", "2. Regulatory + 6. AI Readiness"),
    ("Family Office", "3. Protect from Disruption", "5. Scale (no headcount)"),
    ("Firm with EU Exposure", "2. Regulatory (DORA)", "4. Competitive Advantage"),
    ("Post-Incident Firm", "3. Protect from Disruption", "1. Capital (LP recovery)"),
    ("Firm Adopting AI Tools", "6. AI Readiness", "2. Regulatory + 3. Protection"),
]

table_shape = slide.shapes.add_table(len(rows) + 1, 3,
                                      Inches(0.6), Inches(2.4),
                                      Inches(12.1), Inches(4.5))
table = table_shape.table

# Header row
for col_idx, header in enumerate(["Client Type", "Lead With", "Support With"]):
    cell = table.cell(0, col_idx)
    cell.text = header
    cell.fill.solid()
    cell.fill.fore_color.rgb = DB_DARK
    for p in cell.text_frame.paragraphs:
        p.font.size = Pt(13)
        p.font.color.rgb = WHITE
        p.font.bold = True
        p.font.name = 'Calibri'

# Data rows
for row_idx, (client, lead, support) in enumerate(rows):
    for col_idx, val in enumerate([client, lead, support]):
        cell = table.cell(row_idx + 1, col_idx)
        cell.text = val
        cell.fill.solid()
        cell.fill.fore_color.rgb = WHITE if row_idx % 2 == 0 else DB_GRAY_LIGHT
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(12)
            p.font.color.rgb = DARK_TEXT
            p.font.name = 'Calibri'
            if col_idx == 0:
                p.font.bold = True

# Column widths
table.columns[0].width = Inches(3.6)
table.columns[1].width = Inches(4.2)
table.columns[2].width = Inches(4.3)


# ── SLIDE 10: Measuring Outcome Delivery ────────────────────────────

slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
dark_header_bar(slide)

add_text_box(slide, Inches(0.6), Inches(1.2), Inches(12), Inches(0.6),
             "Measuring Outcome Delivery", font_size=32, color=DB_DARK, bold=True)

add_text_box(slide, Inches(0.6), Inches(1.8), Inches(11), Inches(0.5),
             "Each outcome has observable evidence. Gather these proof points for QBRs and renewals.",
             font_size=14, color=MID_TEXT)

metrics = [
    ("1", "Capital Confidence",
     "DDQs completed & avg. turnaround  |  ODD reviews supported  |  LP-facing reports delivered  |  Allocator retention rate"),
    ("2", "Regulatory Readiness",
     "Policies updated & current  |  Exam readiness score  |  Regulatory changes addressed  |  Zero deficiency findings"),
    ("3", "Protection & Resilience",
     "Risk score trend  |  Vulnerabilities remediated  |  IR plan tested (date)  |  Incidents contained  |  Phishing test results"),
    ("4", "Competitive Advantage",
     "Peer benchmarking position  |  New capabilities implemented  |  Frameworks achieved  |  LP feedback quality"),
    ("5", "Scale & Efficiency",
     "Cost vs. in-house equivalent  |  COO/CCO hours saved  |  New entities covered  |  Services consumed per dollar"),
    ("6", "AI Readiness",
     "AI policy in place & current  |  Employees trained (%)  |  AI tools inventoried & assessed  |  Shadow AI identified  |  Governance cadence"),
]

for i, (num, title, evidence) in enumerate(metrics):
    y = Inches(2.4) + Inches(i * 0.82)

    # Number circle
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.6), y + Inches(0.05),
                                     Inches(0.5), Inches(0.5))
    circle.fill.solid()
    circle.fill.fore_color.rgb = DB_AQUA
    circle.line.fill.background()
    tf = circle.text_frame
    p = tf.paragraphs[0]
    p.text = num
    p.font.size = Pt(16)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    # Row background
    add_rect(slide, Inches(1.3), y, Inches(11.4), Inches(0.7), DB_GRAY_LIGHT)

    add_text_box(slide, Inches(1.5), y + Inches(0.02), Inches(3), Inches(0.35),
                 title, font_size=14, color=DB_DARK, bold=True)
    add_text_box(slide, Inches(1.5), y + Inches(0.35), Inches(11), Inches(0.35),
                 evidence, font_size=11, color=MID_TEXT)


# ── SLIDE 11: What's Next ──────────────────────────────────────────

slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DB_DARK)

add_rect(slide, Inches(0.8), Inches(2.0), Inches(1.2), Inches(0.06), DB_ORANGE)

add_text_box(slide, Inches(0.8), Inches(2.3), Inches(10), Inches(0.8),
             "What's Next", font_size=40, color=WHITE, bold=True)

items = [
    ("Start using the framework now.",
     "Pick one client this week and reframe your next conversation around outcomes instead of deliverables."),
    ("Value narratives are coming.",
     "Detailed stories, data points, and proof that bring each outcome to life in client conversations."),
    ("SKU mapping is coming.",
     "Every service mapped to outcomes so proposals read as 'here's how we deliver Outcome X' — not a task list."),
    ("This is how we become indispensable.",
     "When we're deeply embedded in their compliance, their DDQs, their IR plan, and their board reporting — we're not a vendor anymore. We're infrastructure."),
]

for i, (bold_text, detail_text) in enumerate(items):
    y = Inches(3.3) + Inches(i * 0.85)
    # Bullet accent
    add_rect(slide, Inches(0.8), y + Inches(0.08), Inches(0.15), Inches(0.15), DB_AQUA)
    add_text_box(slide, Inches(1.2), y - Inches(0.03), Inches(10), Inches(0.3),
                 bold_text, font_size=16, color=WHITE, bold=True)
    add_text_box(slide, Inches(1.2), y + Inches(0.3), Inches(10), Inches(0.35),
                 detail_text, font_size=12, color=DB_GRAY)

add_text_box(slide, Inches(0.8), Inches(6.6), Inches(5), Inches(0.4),
             "DRAWBRIDGE  |  CONFIDENTIAL", font_size=11, color=LIGHT_TEXT)


# ── Save ────────────────────────────────────────────────────────────

output_path = "/Users/chrisdelisle/Cursor/Renewal_Forecast/renewal-forecast/Drawbridge_Outcome_Framework.pptx"
prs.save(output_path)
print(f"Saved: {output_path}")
